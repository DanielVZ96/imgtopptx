from pptx import Presentation
from os import *

cpath = input("please enter the full path to the directory you want to work with: ")
chdir(cpath)
foldersinpath = []
supportedexts = ["jpg", "png"]

#Presentation instanciation
newppt = Presentation()

#Layouts that will be used (powerpoint has these slice_layouts indexed):
picture_slide_layout = newppt.slide_layouts[8]
title_slide_layout = newppt.slide_layouts[0]
#Title slide creation
firstslide = newppt.slides.add_slide(title_slide_layout)
title = firstslide.shapes.title
title.text = str(input("Please input a title for your pptx: "))

#Checks every file in the folder selected and makes slides for the images
for file in listdir():
    try:
        filename, extension = file.split(".")
    except:
        foldersinpath.append(file)
        print("{} is a folder and will be iterated over later.".format(file))
        continue
    if extension in supportedexts:
        print("I made a new image slide!")
        fileslide = newppt.slides.add_slide(picture_slide_layout)
        placeholder = fileslide.placeholders[1]
        newimg = placeholder.insert_picture(file)

#Checks for images in the folders inside the path selected
for folder in foldersinpath:
    chdir(folder)
    for file in listdir():
            try:
                filename, extension = file.split(".")
            except:
                foldersinpath.append(file)
                print("{} is a folder and will be iterated over later.".format(file))
                continue
            if extension in supportedexts:
                fileslide = newppt.slides.add_slide(picture_slide_layout)
                placeholder = fileslide.placeholders[1]
                newimg = placeholder.insert_picture(file)
    chdir(cpath)

#Done
newppt.save(input("please enter your desired filename: ") + ".pptx")
